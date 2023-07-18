# -*- coding: utf-8 -*-
"""
Spyder Editor

This is a temporary script file.
"""

from pptx import Presentation
from pptx.util import Cm
import glob
import os
import numpy as np
import tkinter as tk
from tkinter import filedialog

#Input data

#(1)Set the name of your experiment
Inmunos_name = 'test'

#(2)Set the size that you want for your images
img_size = Cm(6)

#(3)Set the size of the step that you want between images putted side by side
space_images = Cm(0.2)

#(4)Set the number of images that you want to plot by slide
n_images = 4

#(5)Set the size of the margen that you want in your ppt
margen_size = Cm(3)

#(6)Set the extension of your files
extension = '.tif'

#Since this point don't modify the script

#########################################

def select_directory():
    # Create a new Tkinter root window for the directory selection
    root = tk.Tk()
    root.withdraw()

    # Prompt the user to select a directory
    directory_path = filedialog.askdirectory()

    # Close the root window to avoid conflicts
    root.destroy()

    # Display the selected directory path
    print(directory_path)
    
    return directory_path
    
# Call the function to select a directory
working_directory = select_directory()
os.chdir(working_directory)

Inmunos_name = working_directory.split('/')
Inmunos_name = Inmunos_name[(len(Inmunos_name)-2)]

# os.chdir('/Volumes/Public/Confocal-June-2023_AM/160623-Cardio-Sections/Combo1/C250622/Processed')

#########################################

########################################
#Automatically get all possible channels to plot

channels_list_full = sorted(glob.glob('*'+extension))

def get_channel_list_from_name(name):
    name_wo_ext = name.replace(extension,'')
    splitted = name_wo_ext.split('_')
    splitted_channel = splitted[(len(splitted)-1)]
    return splitted_channel

channels_list_collapsed = [get_channel_list_from_name(processed_image) for processed_image in channels_list_full]

channels_list_collapsed = list(set(channels_list_collapsed))

#Pop-up window to set the channels you want to plot.

class ItemSelector(tk.Tk):
    def __init__(self, items):
        super().__init__()
        self.items = items
        self.selected_items = []

        self.title("Select channels")

        self.item_listbox = tk.Listbox(self, selectmode=tk.MULTIPLE)
        self.item_listbox.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)

        self.selected_listbox = tk.Listbox(self)
        self.selected_listbox.pack(side=tk.RIGHT, fill=tk.BOTH, expand=True)

        self.move_right_button = tk.Button(self, text=">", command=self.move_right)
        self.move_right_button.pack()

        self.move_left_button = tk.Button(self, text="<", command=self.move_left)
        self.move_left_button.pack()

        self.move_up_button = tk.Button(self, text="Up", command=self.move_up)
        self.move_up_button.pack()

        self.move_down_button = tk.Button(self, text="Down", command=self.move_down)
        self.move_down_button.pack()

        self.update_item_listbox()

    def update_item_listbox(self):
        self.item_listbox.delete(0, tk.END)
        for item in self.items:
            self.item_listbox.insert(tk.END, item)

    def move_right(self):
        selected_indices = self.item_listbox.curselection()
        for index in selected_indices:
            item = self.item_listbox.get(index)
            self.selected_items.append(item)
            self.selected_listbox.insert(tk.END, item)

        self.update_item_listbox()

    def move_left(self):
        selected_indices = self.selected_listbox.curselection()
        for index in selected_indices:
            item = self.selected_listbox.get(index)
            self.selected_items.remove(item)
            self.selected_listbox.delete(index)

        self.update_item_listbox()

    def move_up(self):
        selected_indices = self.selected_listbox.curselection()
        for index in selected_indices:
            if index > 0:
                item = self.selected_listbox.get(index)
                self.selected_items.remove(item)
                self.selected_items.insert(index - 1, item)
                self.selected_listbox.delete(index)
                self.selected_listbox.insert(index - 1, item)

    def move_down(self):
        selected_indices = self.selected_listbox.curselection()
        for index in reversed(selected_indices):
            if index < self.selected_listbox.size() - 1:
                item = self.selected_listbox.get(index)
                self.selected_items.remove(item)
                self.selected_items.insert(index + 1, item)
                self.selected_listbox.delete(index)
                self.selected_listbox.insert(index + 1, item)

    def get_selected_items(self):
        return self.selected_items

# Example usage
items = channels_list_collapsed

selector = ItemSelector(items)
selector.mainloop()
selected_items = selector.get_selected_items()

print("Selected items:", selected_items)

channels_list = selected_items

#######################################

#Data derived from imput

#(1)Calculate the number of channels

n_channels = len(channels_list) #Calculate the number of channels

#(2)Extracting image names

path_folder = os.getcwd()
image_list_full = sorted(glob.glob('*_'+channels_list[0]+extension))
image_list = [image.replace('_'+channels_list[0]+extension,'') for image in image_list_full]

#(3)Function to define ranges that will serve 
#to define the images from the image_list that will be included in each slide

def create_ranges(List):
    
    List_rangles = []
    
    for j in range(0,len(List)-1):
        
        List_rangles.append(range(List[j],List[j+1]))
    
    return List_rangles    
    

number_list = np.arange(0,len(image_list),n_images).tolist()
number_list.append(len(image_list))
range_list = create_ranges(number_list)

#(4)Calculate the highest dimension of your images to know the reference
#image dimension to define steps in coordinates inside a slide

#4a create a temporal presentation with a sample image inside

prs_temp = Presentation() 
prs_temp.slide_width = n_channels*img_size + (n_channels-1)*space_images + margen_size*2
prs_temp.slide_height = n_images*img_size + n_images * Cm(1) + margen_size*2
blank_slide_layout = prs_temp.slide_layouts[6]
slide_temp = prs_temp.slides.add_slide(blank_slide_layout)
img_sample=slide_temp.shapes.add_picture(image_list_full[0],margen_size,margen_size)

#4(b)Create a function with conditional statements to resize images to the highest dimension:

def resize(img_sample,img_size):

    if img_sample.height == img_sample.width:
        
        img_sample.height = img_size
        img_sample.width = img_size
    
    elif img_sample.height > img_sample.width:
        
        Conversion_factor = img_size / Cm(img_sample.height)
        final_width = Cm(img_sample.width) * Conversion_factor
        img_sample.height = img_size
        img_sample.width = int(final_width)
    
    elif img_sample.width > img_sample.height:
        
        Conversion_factor = img_size / Cm(img_sample.width)
        final_height = Cm(img_sample.height) * Conversion_factor
        img_sample.width = img_size
        img_sample.height = int(final_height)

#4(c)get dimensions and calculate y steps

resize(img_sample,img_size)

y_step = img_sample.height
x_step = img_sample.width


del prs_temp
del slide_temp
del img_sample


#Create final presentation

prs = Presentation() #The presentation is called prs
prs.slide_width = n_channels*x_step + (n_channels-1)*space_images + margen_size*2 #Set the width of the slides
prs.slide_height = n_images*y_step + n_images * Cm(1) + margen_size*2 #Set the Height of the slides
blank_slide_layout = prs.slide_layouts[6] #Set the layout of white empty background for the slides

#Code creating the individual slides

for r in range_list:
    
    slide = prs.slides.add_slide(blank_slide_layout) 

    for i,n in zip(range(0,n_images),r):
    
        #Code for adding text box
        
        xpos_title = margen_size
        
        ypos_title = margen_size + y_step * (i) + Cm(1) * i
        
        txBox = slide.shapes.add_textbox(xpos_title, ypos_title, Cm(1), Cm(1))
        tf = txBox.text_frame
        tf.text = image_list[n]
    
    
        #Code for adding Images. 
        
        for j in range(0,len(channels_list)):
            
            path = path_folder +'/' + image_list[n] +'_'+channels_list[j]+extension #path to image
            
            if j == 0:
            
                x_img = margen_size
            
            else:
                
                x_img = margen_size + x_step*(j)
            
            img=slide.shapes.add_picture(path,x_img+space_images*j,ypos_title+Cm(1))
            
            #Conditional statements to resize images to the highest dimension
            
            resize(img,img_size)

#Save presentation

os.chdir('..')
prs.save(Inmunos_name + '.pptx')