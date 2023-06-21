# -*- coding: utf-8 -*-
"""
Spyder Editor

This is a temporary script file.
"""

from pptx import Presentation
from pptx.util import Cm
import glob
import os
import numpy as np

#Input data

#(1)Set the name of your experiment
Inmunos_name = 'Immunofluorescence_Processing'

#(2)Set the size that you want for your images
img_size = Cm(6)

#(3)Set the size of the step that you want between images putted side by side
space_images = Cm(0.2)

#(4)Set the channels you want to plot
channels_list = ['blue','red', 'green','grays', 'Merged_all_wscale']

#(5)Set the number of images that you want to plot by slide
n_images = 1

#(6)Set the size of the margen that you want in your ppt
margen_size = Cm(3)

#(7)Set the path were your processed images are located
os.chdir('/Users/amarco/Desktop/Immunofluorescence_Processing/Processed')

#(8)Set the extension of your images
extension = '.tif'

#Since this point don't modify the script

#Data derived from imput

#(1)Calculate the number of channels

n_channels = len(channels_list) #Calculate the number of channels

#(2)Extracting image names

path_folder = os.getcwd()
image_list_full = sorted(glob.glob('*_'+channels_list[0]+extension))
image_list = [image.replace('_'+channels_list[0]+extension,'') for image in image_list_full]

#(3)Function to define ranges that will serve 
#to define the images from the image_list that will be included in each slide

def create_ranges(List):
    
    List_rangles = []
    
    for j in range(0,len(List)-1):
        
        List_rangles.append(range(List[j],List[j+1]))
    
    return List_rangles    
    

number_list = np.arange(0,len(image_list),n_images).tolist()
number_list.append(len(image_list))
range_list = create_ranges(number_list)

#(4)Calculate the highest dimension of your images to know the reference
#image dimension to define steps in coordinates inside a slide

#4a create a temporal presentation with a sample image inside

prs_temp = Presentation() 
prs_temp.slide_width = n_channels*img_size + (n_channels-1)*space_images + margen_size*2
prs_temp.slide_height = n_images*img_size + n_images * Cm(1) + margen_size*2
blank_slide_layout = prs_temp.slide_layouts[6]
slide_temp = prs_temp.slides.add_slide(blank_slide_layout)
img_sample=slide_temp.shapes.add_picture(image_list_full[0],margen_size,margen_size)

#4(b)Create a function with conditional statements to resize images to the highest dimension:

def resize(img_sample,img_size):

    if img_sample.height == img_sample.width:
        
        img_sample.height = img_size
        img_sample.width = img_size
    
    elif img_sample.height > img_sample.width:
        
        Conversion_factor = img_size / Cm(img_sample.height)
        final_width = Cm(img_sample.width) * Conversion_factor
        img_sample.height = img_size
        img_sample.width = int(final_width)
    
    elif img_sample.width > img_sample.height:
        
        Conversion_factor = img_size / Cm(img_sample.width)
        final_height = Cm(img_sample.height) * Conversion_factor
        img_sample.width = img_size
        img_sample.height = int(final_height)

#4(c)get dimensions and calculate y steps

resize(img_sample,img_size)

y_step = img_sample.height
x_step = img_sample.width


del prs_temp
del slide_temp
del img_sample


#Create final presentation

prs = Presentation() #The presentation is called prs
prs.slide_width = n_channels*x_step + (n_channels-1)*space_images + margen_size*2 #Set the width of the slides
prs.slide_height = n_images*y_step + n_images * Cm(1) + margen_size*2 #Set the Height of the slides
blank_slide_layout = prs.slide_layouts[6] #Set the layout of white empty background for the slides

#Code creating the individual slides

for r in range_list:
    
    slide = prs.slides.add_slide(blank_slide_layout) 

    for i,n in zip(range(0,n_images),r):
    
        #Code for adding text box
        
        xpos_title = margen_size
        
        ypos_title = margen_size + y_step * (i) + Cm(1) * i
        
        txBox = slide.shapes.add_textbox(xpos_title, ypos_title, Cm(1), Cm(1))
        tf = txBox.text_frame
        tf.text = image_list[n]
    
    
        #Code for adding Images. 
        
        for j in range(0,len(channels_list)):
            
            path = path_folder +'/' + image_list[n] +'_'+channels_list[j]+extension #path to image
            
            if j == 0:
            
                x_img = margen_size
            
            else:
                
                x_img = margen_size + x_step*(j)
            
            img=slide.shapes.add_picture(path,x_img+space_images*j,ypos_title+Cm(1))
            
            #Conditional statements to resize images to the highest dimension
            
            resize(img,img_size)

#Save presentation

os.chdir('..')
prs.save(Inmunos_name + '.pptx')
